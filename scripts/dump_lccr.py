import re
from pathlib import Path
from pptx import Presentation
path = Path(r"c:\Users\denis\CrossDevice\Mihai's S24 Ultra (1)\storage\Android\media\com.whatsapp\WhatsApp\Media\WhatsApp Documents\Rupture du LCCR 2026 VV.pptx")
prs = Presentation(str(path))
lines = []
for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            for raw in shape.text.splitlines():
                t = re.sub(r"\s+", " ", raw.strip())
                if len(t) > 8:
                    lines.append(t)
Path("scripts/lccr_lines.txt").write_text("\n".join(lines), encoding="utf-8")
print(len(lines))
