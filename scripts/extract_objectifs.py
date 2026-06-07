"""Extract objectifs d apprentissage and build MCQ questions."""
from __future__ import annotations
import json, random, re, sys
from pathlib import Path
ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT / "scripts"))
from extract_content import SOURCES, SOURCE_LABELS, sanitize
from mcq_choices import add_choices
from pypdf import PdfReader
from pptx import Presentation
OUT_OBJECTIFS = ROOT / "src" / "data" / "objectifs.json"
OUT_QUESTIONS = ROOT / "src" / "data" / "objectifs_questions.json"
MARKER_RE = re.compile(r"\(([AB][+]?)\)\s*([^()]{20,}?)(?=\s*\([AB][+]?\)|$)")
LINE_START_RE = re.compile(r"^\(([AB][+]?)\)\s*(.+)$", re.I)
VERB_RE = re.compile(r"\b(savoir|citer|reconna|connaitre|connaître|expliquer|decrire|décrire|differencier|identifier|evaluer|évaluer|realiser|mettre|appliquer|comprendre|analyser|discuter|presenter|présenter|definir|définir|interpreter|interpréter|choisir|justifier|planifier|enoncer|énoncer|etablir|établir|effectuer|prescrire|informer|definir|définir)\b", re.I)
LCCR_FALLBACK = [
    {"level": "A", "text": "Connaitre les roles mecaniques du ligament croise cranial chez le chien"},
    {"level": "A", "text": "Identifier les facteurs etiologiques de la rupture du LCCR (obesite, conformation, mode de vie)"},
    {"level": "A", "text": "Reconnaitre les signes cliniques evocateurs d une rupture du LCCR"},
    {"level": "A", "text": "Utiliser le test du tiroir cranial dans le diagnostic de rupture du LCCR"},
    {"level": "A", "text": "Interpreter les signes radiographiques associes a la rupture du LCCR"},
    {"level": "B", "text": "Citer les principes des techniques chirurgicales de traitement (TPLO, TTA, techniques extra-articulaires)"},
    {"level": "B", "text": "Decrire les complications possibles apres traitement chirurgical du LCCR"},
    {"level": "A", "text": "Expliquer l interet du LCCr chez le chien comme modele de la rupture du LCA chez l homme"},
]

def read_source_text(path: Path) -> str:
    if path.suffix.lower() == ".pdf":
        return "\n".join((p.extract_text() or "") for p in PdfReader(str(path)).pages)
    prs = Presentation(str(path))
    return "\n".join(sh.text for s in prs.slides for sh in s.shapes if hasattr(sh, "text") and sh.text.strip())

def norm(text: str) -> str:
    return sanitize(text).lower().replace("'", "'")

def is_valid_objective(text: str) -> bool:
    cleaned = sanitize(text)
    if len(cleaned) < 30:
        return False
    if cleaned.count(" ") < 3:
        return False
    if not VERB_RE.search(cleaned):
        return False
    if re.match(r"^[a-z]{1,3}\b", cleaned):
        return False
    return True

def extract_objectifs(text: str) -> list[dict]:
    compact = re.sub(r"\s+", " ", sanitize(text))
    found: list[dict] = []
    for level, body in MARKER_RE.findall(compact):
        objective = sanitize(body)
        if is_valid_objective(objective):
            found.append({"level": level.upper(), "text": objective})
    if found:
        return dedupe_objectifs(found)
    lines = [sanitize(raw) for raw in text.splitlines() if sanitize(raw)]
    for line in lines:
        match = LINE_START_RE.match(line)
        if not match:
            continue
        level, objective = match.group(1).upper(), sanitize(match.group(2))
        if is_valid_objective(objective):
            found.append({"level": level, "text": objective})
    return dedupe_objectifs(found)

def dedupe_objectifs(items: list[dict]) -> list[dict]:
    seen, unique = set(), []
    for item in items:
        key = norm(item["text"])
        if key in seen:
            continue
        seen.add(key)
        unique.append(item)
    return unique

def build_questions(source: str, objectifs: list[dict], all_objectifs_by_source: dict[str, list[dict]]) -> list[dict]:
    questions: list[dict] = []
    all_texts = [obj["text"] for obj in objectifs]
    levels = sorted({obj.get("level", "A") for obj in objectifs})

    for obj in objectifs:
        objective, level = obj["text"], obj.get("level", "A")
        short = objective if len(objective) <= 140 else objective[:137] + "..."
        level_options = list(dict.fromkeys(levels + ["A", "B", "A+", "B+"]))
        level_distractors = [value for value in level_options if value != level][:3]
        questions.append({
            "question": f"Quel est le niveau de cet objectif d apprentissage : « {short} » ?",
            "answer": level,
            "context": f"Objectif d apprentissage (niveau {level}) · {SOURCE_LABELS[source]} · {objective}",
            "source": source,
            "type": "objectif_apprentissage",
            "_distractors": level_distractors[:3],
        })

    wrong_pool: list[str] = []
    for other_source, other_items in all_objectifs_by_source.items():
        if other_source == source:
            continue
        wrong_pool.extend(item["text"] for item in other_items)

    if len(objectifs) >= 3 and len(wrong_pool) >= 2:
        questions.append({
            "question": f"Selectionnez tous les objectifs d apprentissage du cours « {SOURCE_LABELS[source]} ».",
            "answer": objectifs[0]["text"],
            "context": f"Objectifs d'apprentissage · {SOURCE_LABELS[source]}",
            "source": source,
            "type": "objectif_apprentissage",
            "multiSelect": True,
            "_correct_pool": all_texts,
            "_wrong_pool": wrong_pool,
            "num_correct": 2,
        })

    seen, unique = set(), []
    for item in questions:
        key = norm(item["question"])
        if key in seen:
            continue
        seen.add(key)
        unique.append(item)
    return unique
def main() -> None:
    all_objectifs: dict[str, list[dict]] = {}
    for source, path in SOURCES.items():
        if source == "rupture_lccr":
            objectifs = LCCR_FALLBACK
        elif not path.exists():
            print("Skip missing", source)
            continue
        else:
            objectifs = extract_objectifs(read_source_text(path))
        all_objectifs[source] = objectifs
        print(source, len(objectifs))

    all_questions: list[dict] = []
    for source, objectifs in all_objectifs.items():
        all_questions.extend(build_questions(source, objectifs, all_objectifs))
    all_questions = add_choices(all_questions)
    OUT_OBJECTIFS.write_text(json.dumps(all_objectifs, ensure_ascii=False, indent=2), encoding="utf-8")
    OUT_QUESTIONS.write_text(json.dumps(all_questions, ensure_ascii=False, indent=2), encoding="utf-8")
    print("total objectifs", sum(len(v) for v in all_objectifs.values()))
    print("total questions", len(all_questions))

if __name__ == "__main__":
    main()
