"""Extract high-quality quiz questions from course PDFs/PPTX."""
from __future__ import annotations

import json
import re
import sys
from pathlib import Path

try:
    from pypdf import PdfReader
except ImportError:
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "pypdf", "-q"])
    from pypdf import PdfReader

try:
    from pptx import Presentation
except ImportError:
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])
    from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
SOURCES = {
    "fractures": Path(r"c:\Users\denis\Downloads\ilovepdf_merged.pdf"),
    "tumeurs_osseuses": Path(r"c:\Users\denis\Downloads\tumeurs osseuses 2026 VV.pdf"),
    "rupture_lccr": Path(
        r"c:\Users\denis\CrossDevice\Mihai's S24 Ultra (1)\storage\Android\media\com.whatsapp\WhatsApp\Media\WhatsApp Documents\Rupture du LCCR 2026 VV.pptx"
    ),
}

SOURCE_LABELS = {
    "fractures": "Traitement des fractures",
    "tumeurs_osseuses": "Tumeurs osseuses",
    "rupture_lccr": "Rupture du LCCR",
}

GENERIC_PROMPTS = {
    "objectif", "autres cas", "articulaires", "traitement", "diagnostic",
    "resultats", "complications", "indications", "contre-indications", "enva enva",
    "materiel", "conclusion", "perspectives", "importance", "introduction",
}

FRAGMENT_STARTS = (
    "ou ", "de ", "du ", "des ", "en ", "la ", "le ", "les ", "et ", "a ",
    "(", "- ", "• ", "▪ ", "❖ ", "✓ ", "➢ ",
)

NOISE_LINE = re.compile(
    r"^(?:-- \d+ of \d+ --|EnvA|ENVA|\?|◼|✓|❖|➢|o )$",
    re.I,
)


def sanitize(text: str) -> str:
    text = re.sub(r"\bEnvA\b", "", text, flags=re.I)
    text = re.sub(r"\bENVA\b", "", text, flags=re.I)
    text = re.sub(r"[§ü¨Ø]", "", text)
    text = re.sub(r"^[•❖✓➢▪\-–—❑]\s*", "", text)
    text = re.sub(r"\s+", " ", text).strip(" \t-–—→")
    return text.strip()


def normalize_key(text: str) -> str:
    return (
        sanitize(text)
        .lower()
        .replace("'", "'")
        .replace("’", "'")
    )


def is_noise_line(line: str) -> bool:
    if len(line) < 4:
        return True
    if NOISE_LINE.match(line):
        return True
    if re.fullmatch(r"[A-Z0-9\s/\-–—\.]{3,40}", line):
        return True
    return False


def extract_lines(path: Path) -> list[str]:
    if path.suffix.lower() == ".pdf":
        reader = PdfReader(str(path))
        raw_lines = []
        for page in reader.pages:
            raw_lines.extend((page.extract_text() or "").splitlines())
    else:
        prs = Presentation(str(path))
        raw_lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    raw_lines.extend(shape.text.splitlines())

    lines: list[str] = []
    for raw in raw_lines:
        cleaned = sanitize(raw)
        if cleaned and not is_noise_line(cleaned):
            lines.append(cleaned)
    return lines


def is_valid_prompt(prompt: str) -> bool:
    if not (15 <= len(prompt) <= 110):
        return False
    key = normalize_key(prompt)
    if key in GENERIC_PROMPTS:
        return False
    if any(key.startswith(s) for s in FRAGMENT_STARTS):
        return False
    if prompt.endswith(("(", ",", ";", ":", "→", "…", " ou", " et")):
        return False
    if len(prompt.split()) < 3:
        return False
    if re.search(r"\(\s*$", prompt):
        return False
    if "http" in key or "année universitaire" in key:
        return False
    return True


def is_valid_answer(answer: str) -> bool:
    if not (2 <= len(answer) <= 65):
        return False
    if "→" in answer or "EnvA" in answer or "ENVA" in answer:
        return False
    if answer.count("(") != answer.count(")"):
        return False
    if len(answer.split()) > 10:
        return False
    return True


def extract_arrow_questions(lines: list[str], source: str) -> list[dict]:
    questions: list[dict] = []
    label = SOURCE_LABELS[source]

    for i, line in enumerate(lines):
        pairs: list[tuple[str, str]] = []

        if "→" in line:
            parts = [sanitize(p) for p in line.split("→") if sanitize(p)]
            for left, right in zip(parts, parts[1:]):
                pairs.append((left, right))
        elif i + 1 < len(lines) and lines[i + 1].startswith("→"):
            pairs.append((line, sanitize(lines[i + 1].lstrip("→"))))

        for prompt, answer in pairs:
            if not is_valid_prompt(prompt) or not is_valid_answer(answer):
                continue
            questions.append(
                {
                    "question": f"D'apres le cours « {label} », a quoi correspond : « {prompt} » ?",
                    "answer": answer,
                    "context": f"{prompt} → {answer}",
                    "source": source,
                    "type": "relation",
                }
            )
    return questions


def extract_numeric_questions(lines: list[str], source: str) -> list[dict]:
    questions: list[dict] = []
    label = SOURCE_LABELS[source]

    patterns = [
        (
            re.compile(r"^(.{8,90}?)\s*[:：]\s*(\d+(?:[.,]\d+)?\s*(?:%|/\s*100\s*000))\s*$"),
            "Quelle valeur est indiquee dans le cours « {label} » pour : « {prompt} » ?",
        ),
        (
            re.compile(r"^(.{8,90}?)\s*[:：]\s*(\d+(?:[.,]\d+)?\s*(?:mg|kg|cm|mm|h|jours?|semaines?|mois|ans?)(?:\s*/\s*[\w]+)?)\s*$", re.I),
            "Quelle dose, duree ou mesure est indiquee dans le cours « {label} » pour : « {prompt} » ?",
        ),
    ]

    for line in lines:
        for regex, template in patterns:
            match = regex.match(line)
            if not match:
                continue
            prompt = sanitize(match.group(1))
            answer = sanitize(match.group(2))
            if not is_valid_prompt(prompt) and len(prompt) < 8:
                continue
            if len(prompt) > 90 or len(answer) > 40:
                continue
            if not re.search(r"\d", answer):
                continue
            questions.append(
                {
                    "question": template.format(label=label, prompt=prompt),
                    "answer": answer,
                    "context": line,
                    "source": source,
                    "type": "numeric",
                }
            )
    return questions


def extract_fact_questions(lines: list[str], source: str) -> list[dict]:
    questions: list[dict] = []
    label = SOURCE_LABELS[source]

    percent_re = re.compile(
        r"^(.{0,70}?)(\d+(?:[.,]\d+)?)\s*%\s*(.{5,80}?)$",
        re.I,
    )
    bullet_re = re.compile(r"^[•❖✓➢\-–—]\s*(.{20,140})$")

    for line in lines:
        percent_match = percent_re.match(line)
        if percent_match:
            prefix = sanitize(percent_match.group(1))
            value = sanitize(percent_match.group(2)).replace(",", ".") + " %"
            suffix = sanitize(percent_match.group(3))
            statement = sanitize(f"{prefix}{percent_match.group(2)} % {suffix}")
            if len(statement) >= 25 and len(statement) <= 140:
                questions.append(
                    {
                        "question": f"D'apres le cours « {label} », quel pourcentage complete l affirmation : « {statement.replace(value.strip(), '___')} » ?",
                        "answer": value.replace(" %", "%"),
                        "context": line,
                        "source": source,
                        "type": "percent",
                    }
                )
            continue

        bullet_match = bullet_re.match(line)
        if not bullet_match:
            continue
        fact = sanitize(bullet_match.group(1))
        if len(fact) < 25 or len(fact) > 130:
            continue
        if fact.endswith((" ou", " de", " du", " des", " et", "(", ",")):
            continue
        if len(fact.split()) < 5:
            continue

        colon_idx = fact.find(":")
        if 8 <= colon_idx <= 70:
            left = sanitize(fact[:colon_idx])
            right = sanitize(fact[colon_idx + 1 :])
            if is_valid_prompt(left) and is_valid_answer(right):
                questions.append(
                    {
                        "question": f"D'apres le cours « {label} », que signifie ou implique : « {left} » ?",
                        "answer": right,
                        "context": line,
                        "source": source,
                        "type": "definition",
                    }
                )
                continue

        words = fact.split()
        if len(words) < 6:
            continue
        hidden_idx = max(i for i, w in enumerate(words) if len(re.sub(r"[^A-Za-zÀ-ÿ]", "", w)) >= 5)
        hidden = words[hidden_idx].strip(".,;:")
        if len(hidden) < 5 or hidden.startswith("("):
            continue
        blank_words = words.copy()
        blank_words[hidden_idx] = "___"
        blank_sentence = " ".join(blank_words)
        if blank_sentence.count("(") != blank_sentence.count(")"):
            continue
        questions.append(
            {
                "question": f"Complete l affirmation du cours « {label} » : « {blank_sentence} »",
                "answer": hidden,
                "context": line,
                "source": source,
                "type": "fill",
            }
        )

    return questions


def dedupe_questions(questions: list[dict]) -> list[dict]:
    seen: set[tuple[str, str]] = set()
    unique: list[dict] = []
    for q in questions:
        key = (normalize_key(q["question"]), normalize_key(q["answer"]))
        if key in seen:
            continue
        seen.add(key)
        unique.append(q)
    return unique


def main() -> None:
    all_questions: list[dict] = []

    for source_key, path in SOURCES.items():
        if not path.exists():
            print("WARNING missing", path)
            continue

        print("Processing", source_key)
        lines = extract_lines(path)
        print(" ", len(lines), "lines")

        all_questions.extend(extract_arrow_questions(lines, source_key))
        all_questions.extend(extract_numeric_questions(lines, source_key))
        all_questions.extend(extract_fact_questions(lines, source_key))

    all_questions = dedupe_questions(all_questions)

    out_path = ROOT / "src" / "data" / "questions.json"
    out_path.write_text(json.dumps(all_questions, ensure_ascii=False, indent=2), encoding="utf-8")

    by_type: dict[str, int] = {}
    for q in all_questions:
        by_type[q["type"]] = by_type.get(q["type"], 0) + 1

    print("Generated", len(all_questions), "questions")
    print("By type:", by_type)


if __name__ == "__main__":
    main()
